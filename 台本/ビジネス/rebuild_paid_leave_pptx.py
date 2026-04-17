# -*- coding: utf-8 -*-
"""The_Paid_Leave_Paradox.pptx を修正文面で再生成（元は画像スライドのためテキスト非編集）。"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "The_Paid_Leave_Paradox.pptx"
BACKUP = Path(__file__).resolve().parent / "The_Paid_Leave_Paradox_backup_image_only.pptx"

NAVY = RGBColor(0x1A, 0x27, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)
BLUE_ACCENT = RGBColor(0x2E, 0x6B, 0x9E)


def add_textbox(slide, left, top, width, height, text, *, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "メイリオ"
    p.alignment = align
    return box


def set_bg(slide, rgb: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def main():
    if OUT.exists() and not BACKUP.exists():
        import shutil

        shutil.copy2(OUT, BACKUP)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1 表紙 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_textbox(s, 0.8, 2.0, 11.5, 1.2, "実態の声から読み解く 「有給休暇取得」のパラドックス", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        0.8,
        3.4,
        11.5,
        1.4,
        "現代日本の組織における 「価値観の分断」と「調整スキル」の欠如",
        size=18,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(s, 0.8, 6.5, 11.5, 0.5, "組織開発インテリジェンス・レポート", size=12, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

    # --- Slide 2 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.5, "The Sentiment Dashboard", size=11, color=BLUE_ACCENT)
    add_textbox(s, 0.6, 1.0, 12, 1.0, "義務化が進む裏で、\n現場に取り残された「生々しい本音」", size=22, bold=True)
    notes = [
        "【声 A】昭和・平成期には、適宜有給休暇を取得しないと解雇されることもあった。",
        "【声 B】有給……？ 勤続二十数年になるが、取得した記憶にない。",
        "【声 C】付与日数を失効させても、得になるわけではない。",
        "【声 D】周囲が努力しているのに、平然と休むことができるのか。",
    ]
    y = 2.3
    for i, t in enumerate(notes):
        add_textbox(s, 0.7 + (i % 2) * 6.2, y + (i // 2) * 1.35, 5.8, 1.2, t, size=13)
    add_textbox(
        s,
        0.5,
        6.0,
        12.3,
        1.0,
        "制度が整っても、現場の「心理的安全性」と「価値観」は依然として分断されている。",
        size=14,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

    # --- Slide 3 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.4, "合理的視点：未消化の有給休暇は「機会損失」である", size=20, bold=True)
    add_textbox(
        s,
        0.6,
        1.1,
        12,
        1.3,
        "「有給休暇10日を失効させれば、金額にしても10万円以上の損失であり、本来得られるはずだった自由時間も失われる。極めてもったいない。」",
        size=14,
        color=BLUE_ACCENT,
    )
    add_textbox(s, 0.6, 2.6, 5.8, 0.4, "自由時間の消滅", size=14, bold=True)
    add_textbox(s, 0.6, 3.0, 5.8, 0.8, "本来の権利であるはずのリフレッシュ期間の喪失", size=12)
    add_textbox(s, 6.8, 2.6, 5.8, 0.4, "経済的損失", size=14, bold=True)
    add_textbox(s, 6.8, 3.0, 5.8, 0.8, "給与換算で年間10万円以上に相当する労働価値の放棄", size=12)
    bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.0), Inches(13.333), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_textbox(
        s,
        0.5,
        6.25,
        12.3,
        0.8,
        "合理的アプローチ（現代的マインド）から見れば、有給休暇を放棄する行為は極めて非論理的である。",
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # --- Slide 4 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.8, "感情的・伝統的視点：行動を阻む「見えないブレーキ」", size=20, bold=True)
    add_textbox(s, 0.6, 1.3, 11.5, 0.5, "同調圧力", size=14, bold=True, color=BLUE_ACCENT)
    add_textbox(s, 0.7, 1.75, 11.5, 0.9, "「周囲が努力しているにもかかわらず、平然と休むことができるのか」", size=13)
    add_textbox(s, 0.6, 2.7, 11.5, 0.5, "過去のトラウマ", size=14, bold=True, color=BLUE_ACCENT)
    add_textbox(
        s,
        0.7,
        3.15,
        11.5,
        1.1,
        "「昭和・平成期には、適宜取得しなければ解雇される事例もあった。その感覚が残り、取得をためらう者もいる」",
        size=13,
    )
    add_textbox(s, 0.6, 4.35, 11.5, 0.5, "組織への過剰配慮", size=14, bold=True, color=BLUE_ACCENT)
    add_textbox(s, 0.7, 4.8, 11.5, 0.9, "「有給休暇を取得すると、同僚や取引先に負担をかけるのではないか」", size=13)
    b2 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(5.85), Inches(13.333), Inches(1.65))
    b2.fill.solid()
    b2.fill.fore_color.rgb = NAVY
    b2.line.fill.background()
    add_textbox(
        s,
        0.5,
        6.05,
        12.3,
        1.2,
        "insight: 単なる「ためらい」ではなく、過去の労働環境が植え付けた「生存戦略」の名残りがブレーキとなっている。",
        size=12,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # --- Slide 5 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.9, "現場の極端な現実と、企業の「苦肉の策」", size=20, bold=True)
    add_textbox(s, 0.6, 1.2, 5.8, 0.4, "無法地帯（法令形骸化）", size=13, bold=True)
    add_textbox(
        s,
        0.6,
        1.55,
        5.8,
        1.8,
        "「労働基準法が適用されていない」と説明された。\n新型コロナウイルス感染時を除き、日曜・祝日以外に休んだことがない。",
        size=12,
    )
    add_textbox(s, 0.6, 3.5, 5.8, 0.8, "法令が形骸化し、気合と根性のみが支配するレガシー環境。", size=11, color=RGBColor(0x66, 0x66, 0x66))
    add_textbox(s, 6.8, 1.2, 5.8, 0.4, "強制適合ゾーン（計画的付与・強制連休）", size=13, bold=True)
    add_textbox(
        s,
        6.8,
        1.55,
        5.8,
        1.8,
        "「多く取得されると業務に支障があるため、年5日分の取得のみ指示された」。\n「会社指示で4月30日〜5月2日を有給で取得し、その前後は通常勤務とされた」。",
        size=12,
    )
    add_textbox(
        s,
        6.8,
        3.5,
        5.8,
        0.9,
        "業務調整の根本解決を避け、トップダウンで「休ませるアリバイ」を作る防衛策。",
        size=11,
        color=RGBColor(0x66, 0x66, 0x66),
    )

    # --- Slide 6 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.6, "【比較マトリクス】職場に摩擦を生む「価値観の分断」", size=18, bold=True)
    headers = ["観点", "合理的・現代的マインド", "感情的・昭和期の価値観"]
    rows = [
        (
            "休むことへの評価",
            "権利の行使として合理的。「給与を得ながら休めるのに取得しないのは非合理的である」",
            "怠慢や周囲への配慮欠如とみなされやすい。「周囲が努力しているのに、平然と休むことができるのか」",
        ),
        ("労働の動機づけ", "契約と対価（損益の見極め）", "組織への忠誠と自己犠牲"),
        ("会社からの視線", "義務を果たせば自由に休める", "休むことで評価が下がる（解雇への不安の残滓）"),
    ]
    y0 = 1.3
    for ri, row in enumerate(rows):
        y = y0 + ri * 1.55
        add_textbox(s, 0.5, y, 2.5, 1.3, row[0], size=12, bold=True)
        add_textbox(s, 3.1, y, 4.8, 1.45, row[1], size=11)
        add_textbox(s, 8.0, y, 4.8, 1.45, row[2], size=11)

    # --- Slide 7 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_textbox(s, 0.6, 0.5, 12, 0.6, "パラダイムシフト：それは「価値観」ではなく「スキル」の問題である", size=20, bold=True, color=WHITE)
    mid = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(2.0), Inches(2.0), Inches(9.3), Inches(1.6))
    mid.fill.solid()
    mid.fill.fore_color.rgb = WHITE
    mid.line.color.rgb = WHITE
    add_textbox(
        s,
        2.2,
        2.35,
        8.9,
        1.2,
        "休めない要因は、当事者の選択や覚悟の問題ではなく、\n業務調整における能力不足の表れである場合がある。",
        size=20,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        0.7,
        4.2,
        12,
        2.0,
        "有給休暇の取得を妨げている本質は、「休む勇気」や「職場の空気」単体では説明しきれない。\n自身が不在でも業務が継続するよう段取りを組む「業務調整ケイパビリティ（能力）」の不足が、取得を阻む要因となりうる。",
        size=15,
        color=WHITE,
    )

    # --- Slide 8 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.7, "【総合整理】有給休暇取得における「3つの人材タイプ」", size=18, bold=True)
    add_textbox(
        s,
        0.6,
        1.1,
        12,
        0.7,
        "調整して休む計画を立てられる者、調整に課題がある者、調整を経ずに休む者が混在する。",
        size=13,
    )
    add_textbox(s, 0.6, 2.0, 5.8, 0.4, "自律型", size=14, bold=True, color=NAVY)
    add_textbox(
        s,
        0.6,
        2.4,
        5.8,
        2.2,
        "調整して休む計画を立てられる。業務の属人化を防ぎ、周囲と協調しながら計画的に権利を行使する。",
        size=12,
    )
    add_textbox(s, 6.8, 2.0, 5.8, 0.4, "依存型", size=14, bold=True, color=BLUE_ACCENT)
    add_textbox(
        s,
        6.8,
        2.4,
        5.8,
        2.2,
        "調整に課題がある。取得をためらう背景に、業務の属人化や調整コストへの回避が絡む場合がある。",
        size=12,
    )
    add_textbox(s, 0.6, 4.8, 12, 0.4, "破壊型", size=14, bold=True, color=RGBColor(0xB0, 0x44, 0x44))
    add_textbox(
        s,
        0.6,
        5.2,
        12,
        1.5,
        "調整を経ずに休む。周囲への影響を十分に考慮せず、突発的・一方的に休むと、チーム運営に大きな負荷を与えうる。",
        size=12,
    )

    # --- Slide 9 ---
    s = prs.slides.add_slide(blank)
    set_bg(s, RGBColor(0xF5, 0xF3, 0xEF))
    add_textbox(s, 0.5, 0.4, 12, 0.8, "結論：有給休暇の消化は「福利厚生」ではなく「マネジメント指標」である", size=18, bold=True)
    add_textbox(
        s,
        0.6,
        1.2,
        12,
        0.7,
        "「部下への規範となるよう、計画的に有給休暇を取得してください」",
        size=14,
        color=BLUE_ACCENT,
    )
    add_textbox(s, 0.6, 2.1, 3.8, 0.35, "STEP 1 属人化の排除", size=12, bold=True, color=NAVY)
    add_textbox(
        s,
        0.6,
        2.45,
        3.8,
        1.5,
        "チーム内で「誰が休んでも業務が回る」仕組みを構築する。有給休暇消化率は、「業務の標準化レベル」を示す指標になりうる。",
        size=11,
    )
    add_textbox(s, 4.7, 2.1, 3.8, 0.35, "STEP 2 リーダーの率先垂範", size=12, bold=True, color=NAVY)
    add_textbox(
        s,
        4.7,
        2.45,
        3.8,
        1.5,
        "マネジメント層が「調整して休む（自律型）」姿を示し、過度な同調圧力を緩和する。",
        size=11,
    )
    add_textbox(s, 8.8, 2.1, 3.8, 0.35, "STEP 3 評価軸の転換", size=12, bold=True, color=NAVY)
    add_textbox(
        s,
        8.8,
        2.45,
        3.8,
        1.5,
        "「休まず働くこと」より、「休むための調整力」を高く評価する文化を育てる。",
        size=11,
    )
    b3 = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(5.85), Inches(13.333), Inches(1.65))
    b3.fill.solid()
    b3.fill.fore_color.rgb = NAVY
    b3.line.fill.background()
    add_textbox(
        s,
        0.6,
        6.05,
        12.2,
        1.3,
        "「計画的な有給休暇の取得が、組織と個人の双方にとって望ましい」という認識を共有することが、現代のリーダーシップである。",
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    prs.save(OUT)
    print("Saved:", OUT)
    if BACKUP.exists():
        print("Backup:", BACKUP)


if __name__ == "__main__":
    main()
